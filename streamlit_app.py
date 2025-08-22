import streamlit as st
from PIL import Image
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches

# Configuração da página
st.set_page_config(page_title="Análise de Publicidade NWE", layout="wide")
st.title("📊 Gerador de Apresentações de Publicidade - NWE")
st.markdown("Faça upload das imagens das campanhas publicitárias e gere automaticamente uma apresentação com gráficos e insights.")

# Upload de arquivos
uploaded_files = st.file_uploader(
    "📁 Envie as imagens das campanhas",
    type=["png", "jpg", "jpeg"],
    accept_multiple_files=True
)

if uploaded_files:
    st.success(f"{len(uploaded_files)} imagem(ns) carregada(s). Clique abaixo para gerar a apresentação.")

    if st.button("🎯 Gerar Apresentação"):
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[5]

        # Slide de título
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "Análise de Publicidade - Loja NWE"
        slide.placeholders[1].text = "Relatório gerado automaticamente com base nas imagens enviadas."

        # Slides com imagens
        temp_dir = tempfile.gettempdir()
        for img_file in uploaded_files:
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = f"Campanha: {img_file.name}"

            image_path = os.path.join(temp_dir, img_file.name)
            with open(image_path, "wb") as f:
                f.write(img_file.getbuffer())

            left = Inches(1)
            top = Inches(2)
            height = Inches(4.5)
            slide.shapes.add_picture(image_path, left, top, height=height)

        # Salvar apresentação
        output_path = "Apresentacao_Publicidade_NWE.pptx"
        prs.save(output_path)
        st.success("✅ Apresentação gerada com sucesso!")

        # Botão para download
        with open(output_path, "rb") as f:
            st.download_button("📥 Baixar Apresentação", f, file_name=output_path)


